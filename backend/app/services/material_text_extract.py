"""
Extract plain text from uploaded learning-material bytes for quiz generation.
"""

from __future__ import annotations

import io
import zipfile
import xml.etree.ElementTree as ET
from typing import Optional, Tuple

# XML namespace for DrawingML text (used in PPTX slides)
_DRAWINGML_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _extract_pptx_xml_fallback(data: bytes) -> str:
    """
    Fallback: extract text directly from PPTX XML without python-pptx.

    Reads every ppt/slides/slideN.xml and ppt/notesSlides/notesSlideN.xml
    entry and collects all <a:t> text nodes.  Works even when python-pptx
    cannot fully parse a shape type.
    """
    texts: list[str] = []
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as zf:
            names = zf.namelist()
            slide_paths = sorted(
                n for n in names
                if n.startswith("ppt/slides/slide") and n.endswith(".xml")
            )
            notes_paths = sorted(
                n for n in names
                if n.startswith("ppt/notesSlides/") and n.endswith(".xml")
            )
            for path in slide_paths + notes_paths:
                try:
                    xml_bytes = zf.read(path)
                    root = ET.fromstring(xml_bytes)
                    slide_texts: list[str] = []
                    for node in root.iter(f"{{{_DRAWINGML_NS}}}t"):
                        t = (node.text or "").strip()
                        if t:
                            slide_texts.append(t)
                    if slide_texts:
                        texts.append(" ".join(slide_texts))
                except Exception:
                    pass
    except Exception:
        pass
    return "\n\n".join(texts)


def extract_pptx_text(data: bytes) -> str:
    """
    Extract all readable text from a PPTX file.

    Covers:
    - All text frames and placeholders (paragraph/run level)
    - Tables (row → cell text)
    - Group shapes (recursive iteration)
    - Speaker notes
    - Falls back to raw XML parsing if python-pptx gives too little text
    """
    from pptx import Presentation

    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        _GROUP = MSO_SHAPE_TYPE.GROUP
    except Exception:
        _GROUP = None

    def _iter_shapes(shapes):
        """Recursively yield text lines from all shapes."""
        for shape in shapes:
            # Recurse into group shapes
            try:
                if _GROUP is not None and shape.shape_type == _GROUP:
                    if hasattr(shape, "shapes"):
                        yield from _iter_shapes(shape.shapes)
                    continue
            except Exception:
                pass

            # Tables: iterate rows → cells
            try:
                if shape.has_table:
                    for row in shape.table.rows:
                        row_parts: list[str] = []
                        for cell in row.cells:
                            try:
                                cell_text = (cell.text_frame.text or "").strip()
                            except Exception:
                                cell_text = ""
                            if cell_text:
                                row_parts.append(cell_text)
                        if row_parts:
                            yield " | ".join(row_parts)
                    continue
            except Exception:
                pass

            # Text frames: paragraph → run level extraction
            try:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        # Prefer joining individual runs (preserves spacing better)
                        runs_text = "".join(r.text for r in para.runs if r.text).strip()
                        if not runs_text:
                            runs_text = (para.text or "").strip()
                        if runs_text:
                            yield runs_text
                    continue
            except Exception:
                pass

            # Last-resort: shape.text property
            try:
                t = (shape.text or "").strip()
                if t:
                    yield t
            except Exception:
                pass

    try:
        presentation = Presentation(io.BytesIO(data))
    except Exception as exc:
        # python-pptx cannot open file — go straight to XML fallback
        fallback = _extract_pptx_xml_fallback(data)
        if fallback:
            return fallback
        raise exc

    slide_chunks: list[str] = []
    for slide_num, slide in enumerate(presentation.slides, 1):
        slide_lines = list(_iter_shapes(slide.shapes))

        # Include speaker notes (often contain detailed explanations)
        try:
            if slide.has_notes_slide:
                notes_text = (
                    slide.notes_slide.notes_text_frame.text or ""
                ).strip()
                if notes_text:
                    slide_lines.append(f"[Notes] {notes_text}")
        except Exception:
            pass

        if slide_lines:
            slide_chunks.append(
                f"[Slide {slide_num}]\n" + "\n".join(slide_lines)
            )

    result = "\n\n".join(slide_chunks)

    # XML fallback: if python-pptx gave very little text, try raw XML parsing
    if len(result.strip()) < 500:
        xml_result = _extract_pptx_xml_fallback(data)
        if len(xml_result.strip()) > len(result.strip()):
            result = xml_result

    return result


def extract_pdf_text(data: bytes) -> str:
    """
    Extract text from PDF pages.

    Tries PyPDF2 first; concatenates pages with double-newline separators
    to preserve page structure.  Returns empty string on failure rather
    than raising.
    """
    import PyPDF2

    try:
        reader = PyPDF2.PdfReader(io.BytesIO(data))
        pages: list[str] = []
        for page in reader.pages:
            try:
                text = (page.extract_text() or "").strip()
            except Exception:
                text = ""
            if text:
                pages.append(text)
        return "\n\n".join(pages)
    except Exception:
        return ""


def extract_docx_text(data: bytes) -> str:
    try:
        import docx  # python-docx

        document = docx.Document(io.BytesIO(data))
        return "\n".join(p.text for p in document.paragraphs if p.text)
    except ImportError:
        # Minimal fallback: read word/document.xml from the docx zip.
        with zipfile.ZipFile(io.BytesIO(data)) as zf:
            xml = zf.read("word/document.xml")
        root = ET.fromstring(xml)
        texts = []
        for node in root.iter():
            if node.tag.endswith("}t") and node.text:
                texts.append(node.text)
        return " ".join(texts)


def _guess_kind(file_type: str, filename: str, content_type: str = "") -> str:
    ft = (file_type or "").lower().strip()
    if ft in {"pdf", "pptx", "ppt", "docx", "doc", "txt", "text"}:
        return ft
    name = (filename or "").lower()
    for ext in ("pdf", "pptx", "ppt", "docx", "doc", "txt"):
        if name.endswith(f".{ext}"):
            return ext
    ctype = (content_type or "").lower()
    if "pdf" in ctype:
        return "pdf"
    if "presentation" in ctype or "powerpoint" in ctype:
        return "pptx"
    if "word" in ctype:
        return "docx"
    if "text/plain" in ctype:
        return "txt"
    return ft or "unknown"


def extract_text_from_bytes(
    data: bytes,
    *,
    file_type: str = "",
    filename: str = "",
    content_type: str = "",
) -> Tuple[str, Optional[str]]:
    """
    Return (extracted_text, error_message).

    error_message is set when extraction fails or the file type is unsupported.
    """
    if not data:
        return "", "empty_file"

    kind = _guess_kind(file_type, filename, content_type)
    try:
        if kind == "pdf":
            return extract_pdf_text(data), None
        if kind in {"txt", "text"}:
            return data.decode("utf-8", errors="ignore"), None
        if kind in {"pptx", "ppt"}:
            return extract_pptx_text(data), None
        if kind in {"docx", "doc"}:
            return extract_docx_text(data), None
        return "", f"unsupported_file_type:{kind or 'unknown'}"
    except Exception as exc:
        return "", f"extraction_failed:{exc}"
